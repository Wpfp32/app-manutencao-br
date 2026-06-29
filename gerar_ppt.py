from pptx import Presentation

# Cria a apresentação
prs = Presentation()

# --- Slide 1: Capa ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Plataforma Digital de Gestão de Manutenção de Frota"
slide.placeholders[1].text = "BR Construções: Inteligência de Dados e Previsibilidade Operacional\nApresentador: Wanderson Phillype"

# --- Slide 2: O Desafio ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Onde estamos perdendo eficiência?"
tf = slide.placeholders[1].text_frame
tf.text = "Controle Manual e Fragmentado: Planilhas e papel geram atrasos."
tf.add_paragraph().text = "Revisões Baseadas em Suposições: Falta de visibilidade do desgaste real."
tf.add_paragraph().text = "Manutenção Corretiva vs. Preventiva: Apagar incêndios custa caro."
tf.add_paragraph().text = "Falta de Mobilidade: O dado do canteiro demora a chegar na gestão."

# --- Slide 3: A Solução ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "A Solução Desenvolvida"
tf = slide.placeholders[1].text_frame
tf.text = "Sistema Web/Mobile Responsivo: Acesso no campo e escritório."
tf.add_paragraph().text = "Banco de Dados em Nuvem (PostgreSQL): Histórico permanente e seguro."
tf.add_paragraph().text = "Controle de Acesso: Perfis de Operador e Gestor."
tf.add_paragraph().text = "Operação Offline-Ready: Upload em lote para áreas remotas."

# --- Slide 4: Dashboard ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Visão Estratégica e Dashboard"
tf = slide.placeholders[1].text_frame
tf.text = "Status Instantâneo: Normal, Atenção e Vencida (Semáforo)."
tf.add_paragraph().text = "Previsão de Parada: Algoritmo calcula a data da próxima manutenção."
tf.add_paragraph().text = "Histórico de Execução: Registro imutável de todas as revisões."

# --- Slide 5: Benefícios ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "O que a BR Construções ganha hoje?"
tf = slide.placeholders[1].text_frame
tf.text = "Redução de Custos (Opex): Transição para manutenção preditiva."
tf.add_paragraph().text = "Aumento da Disponibilidade: Menos máquinas paradas no canteiro."
tf.add_paragraph().text = "Fim do Retrabalho Administrativo: Atualização instantânea."
tf.add_paragraph().text = "Auditoria e Transparência: Ciclo de vida documentado."

# --- Slide 6: Próximos Passos ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Próximos Passos"
tf = slide.placeholders[1].text_frame
tf.text = "Fase 1: Cadastro completo do inventário atual da frota."
tf.add_paragraph().text = "Fase 2: Treinamento de 15 minutos com os operadores."
tf.add_paragraph().text = "Fase 3: Acompanhamento do primeiro ciclo preditivo."
tf.add_paragraph().text = "\nDúvidas e Sugestões?"

# Salva o arquivo na sua pasta
prs.save('Apresentacao_BR_Construcoes.pptx')
print("✅ Arquivo 'Apresentacao_BR_Construcoes.pptx' gerado com sucesso na sua pasta!")